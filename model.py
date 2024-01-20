import os
import smtplib
import time
from email.mime.application import MIMEApplication
from email.mime.multipart import MIMEMultipart
from email.mime.text import MIMEText

import fitz  # PyMuPDF
import pytesseract
from PIL import Image
from pptx import Presentation
from pptx.util import Pt


def getImgsFromPDF(pdf_path, output_folder):
    os.makedirs(output_folder, exist_ok=True)

    pdf_document = fitz.open(pdf_path)

    for page_index in range(len(pdf_document)):
        page = pdf_document[page_index]

        image_list = page.get_images(full=True)

        if image_list:
            print(f"[+] Found a total of {len(image_list)} images in page {page_index}")
        else:
            print("[!] No images found on page", page_index)

        for image_index, img in enumerate(image_list):
            xref = img[0]

            base_image = pdf_document.extract_image(xref)
            image_bytes = base_image["image"]

            image_ext = base_image["ext"]

            image_name = f"image{page_index + 1}_{image_index + 1}.{image_ext}"

            image_path = os.path.join(output_folder, image_name)
            with open(image_path, "wb") as image_file:
                image_file.write(image_bytes)

            print(f"Saved image {image_index + 1} from page {page_index + 1}")

    pdf_document.close()


def getTextFromImgs(images_folder, text_folder):
    """
    Extracts text from all the images in images_folder and saves the resulting text files in text_folder.
    """
    # Ensure output directory exists
    os.makedirs(text_folder, exist_ok=True)

    for image_name in os.listdir(images_folder):
        if image_name.endswith((".png", ".jpg", ".jpeg", ".tiff", ".bmp", ".gif")):
            image_path = os.path.join(images_folder, image_name)

            image = Image.open(image_path)

            text = pytesseract.image_to_string(image, lang='eng')

            text_file_name = os.path.splitext(image_name)[0] + '.txt'
            text_file_path = os.path.join(text_folder, text_file_name)

            with open(text_file_path, 'w', encoding='utf-8') as text_file:
                text_file.write(text)

            print(f"Processed {image_name}")


def getPPTFromImgText(images_folder, ppt_path):
    """
    Creates a PowerPoint presentation with the text extracted from images in images_folder and saves it to ppt_path.
    """
    prs = Presentation()
    for image_name in sorted(os.listdir(images_folder)):
        if image_name.endswith((".png", ".jpg", ".jpeg", ".tiff", ".bmp", ".gif")):
            image_path = os.path.join(images_folder, image_name)

            image = Image.open(image_path)
            text = pytesseract.image_to_string(image, lang='eng')

            slide_layout = prs.slide_layouts[5]  # choosing a blank slide layout
            slide = prs.slides.add_slide(slide_layout)

            txBox = slide.shapes.add_textbox(Pt(50), Pt(50), prs.slide_width - Pt(100), prs.slide_height - Pt(100))
            tf = txBox.text_frame
            tf.text = text  # Add the extracted text to the text frame

    prs.save(ppt_path)
    print(f"PowerPoint saved at {ppt_path}")


def send_email_with_ppt(recipient_email, pptx_file_path):
    time.sleep(3)

    HOST = "smtp-mail.outlook.com"
    PORT = 587
    FROM_EMAIL = "codesteinsprojectmail@gmail.com"
    TO_EMAIL = recipient_email
    PASSWORD = "codes#2oo345"

    msg = MIMEMultipart()
    msg['From'] = FROM_EMAIL
    msg['To'] = TO_EMAIL
    msg['Subject'] = "Here is the required Powerpoint Presentation"

    message = "Please find the attached Powerpoint presentation."
    msg.attach(MIMEText(message, 'plain'))

    with open(pptx_file_path, 'rb') as file:
        pptx_attachment = MIMEApplication(file.read(), _subtype="pptx")
    pptx_attachment.add_header('Content-Disposition', f'attachment; filename={pptx_file_path}')
    msg.attach(pptx_attachment)

    try:
        smtp = smtplib.SMTP(HOST, PORT)
        smtp.starttls()
        smtp.login(FROM_EMAIL, PASSWORD)
        smtp.sendmail(FROM_EMAIL, TO_EMAIL, msg.as_string())
        smtp.quit()
        print("Email sent successfully!")

    except Exception as e:
        print(f"Email could not be sent. Error: {e}")


def send_feedback_email(recipient_email, feedback_text):
    HOST = "smtp-mail.outlook.com"  # or your SMTP server
    PORT = 587
    FROM_EMAIL = "codesteinsprojectmail@gmail.com"  # your email
    PASSWORD = "codes#2oo345"  # your email password

    msg = MIMEMultipart('alternative')
    msg['From'] = FROM_EMAIL
    msg['To'] = recipient_email
    msg['Subject'] = "New Feedback Submission"

    body_content = """
    <html>
    <head>
        <style>
            body {{
                font-family: Arial, sans-serif;
                margin: 20px;
                background-color: #f4f4f4;
                color: #555;
            }}
            .container {{
                background-color: #fff;
                padding: 20px;
                border-radius: 5px;
                box-shadow: 0 0 10px rgba(0, 0, 0, 0.1);
            }}
            .header {{
                font-size: 24px;
                font-weight: bold;
                margin-bottom: 20px;
            }}
            .feedback-text {{
                margin-top: 20px;
                line-height: 1.6;
            }}
            .footer {{
                margin-top: 30px;
                font-size: 12px;
                text-align: center;
            }}
        </style>
    </head>
    <body>
        <div class="container">
            <div class="header">Feedback Received</div>
            <div class="feedback-text">
                <p><b>Feedback:</b></p>
                <p>{feedback_text}</p>
            </div>
            <div class="footer">
                This is an automated message. Please do not reply directly to this email.
            </div>
        </div>
    </body>
    </html>
    """.format(feedback_text=feedback_text)

    part2 = MIMEText(body_content, 'html')
    msg.attach(part2)

    try:
        smtp = smtplib.SMTP(HOST, PORT)
        smtp.starttls()
        smtp.login(FROM_EMAIL, PASSWORD)
        smtp.sendmail(FROM_EMAIL, recipient_email, msg.as_string())
        smtp.quit()
        print("Feedback email sent successfully!")

    except Exception as e:
        print(f"Feedback email could not be sent. Error: {e}")
